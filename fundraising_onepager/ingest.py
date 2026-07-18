"""Deck and website extraction.

Both extractors return plain text with structural markers (slide/page breaks,
heading prefixes) so the model can reason about deck order and site hierarchy
without us shipping a layout model.
"""

from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

SUPPORTED_DECK_SUFFIXES = {".pdf", ".pptx"}

# Cap what we send to the API. A 60-slide deck is ~30-60k chars; the model has
# a 1M window, so this is about cost and noise, not capacity.
MAX_DECK_CHARS = 200_000
MAX_SITE_CHARS = 40_000

USER_AGENT = (
    "Mozilla/5.0 (compatible; FundraisingReadinessBot/1.0; "
    "+https://example.invalid/bot)"
)
FETCH_TIMEOUT = 20  # seconds


class IngestError(Exception):
    """Raised when a deck or site cannot be turned into usable text."""


@dataclass
class DeckContent:
    text: str
    unit_count: int  # slides for PPTX, pages for PDF
    source: Path

    @property
    def unit_label(self) -> str:
        return "slide" if self.source.suffix.lower() == ".pptx" else "page"


@dataclass
class SiteContent:
    url: str
    title: str
    headings: list[str]
    ctas: list[str]
    body: str

    def as_prompt_text(self) -> str:
        parts = [f"URL: {self.url}", f"TITLE: {self.title or '(none)'}"]
        if self.headings:
            parts.append("HEADINGS:\n" + "\n".join(f"- {h}" for h in self.headings))
        if self.ctas:
            parts.append(
                "VISIBLE CTAS / NAV LINKS:\n" + "\n".join(f"- {c}" for c in self.ctas)
            )
        parts.append("BODY COPY:\n" + self.body)
        return "\n\n".join(parts)


def _collapse(text: str) -> str:
    """Squash runaway whitespace without losing paragraph breaks."""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


# --------------------------------------------------------------------------- deck


def extract_deck(path: Path) -> DeckContent:
    """Extract ordered text from a .pdf or .pptx deck."""
    if not path.exists():
        raise IngestError(f"Deck not found: {path}")
    if not path.is_file():
        raise IngestError(f"Deck path is not a file: {path}")

    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_DECK_SUFFIXES:
        supported = ", ".join(sorted(SUPPORTED_DECK_SUFFIXES))
        raise IngestError(
            f"Unsupported deck type '{suffix or '(none)'}'. Supported: {supported}"
        )

    content = _extract_pdf(path) if suffix == ".pdf" else _extract_pptx(path)

    if len(content.text) < 200:
        raise IngestError(
            f"Extracted only {len(content.text)} characters of text from {path.name}. "
            "The deck may be image-only (scanned or exported as flat images) — "
            "this tool reads text, not pixels. Try a text-based export."
        )
    return content


def _extract_pdf(path: Path) -> DeckContent:
    import pdfplumber

    pages: list[str] = []
    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"--- Page {i} ---\n{text.strip()}")
    except IngestError:
        raise
    except Exception as exc:  # pdfplumber raises a wide variety on bad files
        raise IngestError(f"Could not read PDF {path.name}: {exc}") from exc

    return DeckContent(
        text=_collapse("\n\n".join(pages))[:MAX_DECK_CHARS],
        unit_count=len(pages),
        source=path,
    )


def _extract_pptx(path: Path) -> DeckContent:
    from pptx import Presentation

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise IngestError(f"Could not read PPTX {path.name}: {exc}") from exc

    slides: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        chunks: list[str] = []
        for shape in slide.shapes:
            chunks.extend(_shape_text(shape))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                chunks.append(f"[Speaker notes] {notes}")
        body = "\n".join(c for c in chunks if c.strip())
        if body.strip():
            slides.append(f"--- Slide {i} ---\n{body.strip()}")

    return DeckContent(
        text=_collapse("\n\n".join(slides))[:MAX_DECK_CHARS],
        unit_count=len(slides),
        source=path,
    )


def _shape_text(shape) -> list[str]:
    """Pull text out of a shape, recursing into groups and tables."""
    out: list[str] = []
    if shape.shape_type == 6 and hasattr(shape, "shapes"):  # MSO_SHAPE_TYPE.GROUP
        for child in shape.shapes:
            out.extend(_shape_text(child))
        return out
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                out.append(" | ".join(cells))
        return out
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            out.append(text)
    return out


# --------------------------------------------------------------------------- site

_STRIP_TAGS = ("script", "style", "noscript", "svg", "iframe", "template")
_CHROME_TAGS = ("nav", "header", "footer", "aside", "form")


def fetch_site(url: str) -> SiteContent:
    """Fetch a URL and extract headings, CTAs, and main body copy."""
    import requests
    from bs4 import BeautifulSoup

    if not url.startswith(("http://", "https://")):
        url = "https://" + url

    try:
        resp = requests.get(
            url,
            timeout=FETCH_TIMEOUT,
            headers={"User-Agent": USER_AGENT, "Accept": "text/html,*/*"},
        )
        resp.raise_for_status()
    except requests.exceptions.SSLError as exc:
        raise IngestError(f"TLS error fetching {url}: {exc}") from exc
    except requests.exceptions.Timeout as exc:
        raise IngestError(f"Timed out after {FETCH_TIMEOUT}s fetching {url}") from exc
    except requests.exceptions.HTTPError as exc:
        raise IngestError(
            f"{url} returned HTTP {resp.status_code}. "
            "The site may block automated clients."
        ) from exc
    except requests.exceptions.RequestException as exc:
        raise IngestError(f"Could not reach {url}: {exc}") from exc

    ctype = resp.headers.get("content-type", "")
    if "html" not in ctype.lower():
        raise IngestError(f"{url} returned '{ctype or 'unknown'}', expected HTML.")

    soup = BeautifulSoup(resp.text, "lxml")

    title = (soup.title.get_text(strip=True) if soup.title else "") or ""

    # CTAs come from chrome, so read them before stripping it.
    ctas = _dedupe(
        el.get_text(" ", strip=True)
        for el in soup.select("a, button, [role='button']")
        if 2 <= len(el.get_text(" ", strip=True)) <= 60
    )[:25]

    for tag in soup(list(_STRIP_TAGS)):
        tag.decompose()

    headings = _dedupe(
        h.get_text(" ", strip=True) for h in soup.select("h1, h2, h3")
    )[:30]

    main = soup.find("main") or soup.find(attrs={"role": "main"}) or soup.body or soup
    for tag in main(list(_CHROME_TAGS)):
        tag.decompose()

    body = _collapse(main.get_text("\n", strip=True))[:MAX_SITE_CHARS]

    if len(body) < 100 and not headings:
        raise IngestError(
            f"Extracted almost no readable text from {url}. The page is likely "
            "client-side rendered; this tool does not run JavaScript."
        )

    return SiteContent(url=url, title=title, headings=headings, ctas=ctas, body=body)


def _dedupe(items) -> list[str]:
    """Order-preserving dedupe, case-insensitive, dropping empties."""
    seen: set[str] = set()
    out: list[str] = []
    for item in items:
        key = item.lower()
        if item and key not in seen:
            seen.add(key)
            out.append(item)
    return out
